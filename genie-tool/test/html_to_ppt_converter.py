#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
HTML转PPT转换器

提供多种方案将生成的HTML格式PPT转换为真正的PowerPoint文件(.pptx)
"""

import os
import sys
import asyncio
from typing import Optional, List
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("请安装python-pptx: pip install python-pptx")
    sys.exit(1)

try:
    from bs4 import BeautifulSoup
except ImportError:
    print("请安装beautifulsoup4: pip install beautifulsoup4")
    sys.exit(1)

try:
    import selenium
    from selenium import webdriver
    from selenium.webdriver.chrome.options import Options
    from selenium.webdriver.common.by import By
    import time
except ImportError:
    selenium = None
    print("⚠️  Selenium未安装，将跳过浏览器截图功能")


class HTMLToPPTConverter:
    """HTML转PPT转换器"""
    
    def __init__(self):
        """初始化转换器"""
        self.presentation = None
        
    def method1_parse_html_content(self, html_file: str, output_file: str = None) -> str:
        """
        方案1：解析HTML内容生成PPTX
        
        优点：纯Python实现，不依赖浏览器
        缺点：无法完美还原复杂样式和图表
        """
        print("🔄 使用方案1：HTML内容解析转换...")
        
        if output_file is None:
            output_file = html_file.replace('.html', '.pptx')
        
        # 读取HTML文件
        with open(html_file, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        # 解析HTML
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # 创建PPT
        prs = Presentation()
        
        # 查找所有幻灯片
        slides = soup.find_all('div', class_=['slide', 'ppt-slide'])
        
        if not slides:
            # 如果没有找到标准的slide结构，尝试其他方式
            slides = soup.find_all('section') or [soup]
        
        for i, slide_div in enumerate(slides):
            print(f"  处理第 {i+1} 页...")
            
            # 添加幻灯片
            slide_layout = prs.slide_layouts[5]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 提取标题
            title_elem = (slide_div.find('h1') or 
                         slide_div.find('h2') or 
                         slide_div.find('.title') or
                         slide_div.find('[class*="title"]'))
            
            if title_elem:
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(9), Inches(1)
                )
                title_frame = title_box.text_frame
                title_frame.text = title_elem.get_text().strip()
                title_frame.paragraphs[0].font.size = Pt(32)
                title_frame.paragraphs[0].font.bold = True
            
            # 提取内容
            content_y = 2.0
            for elem in slide_div.find_all(['p', 'ul', 'ol', 'div']):
                if elem.name in ['ul', 'ol']:
                    # 处理列表
                    for li in elem.find_all('li'):
                        text = li.get_text().strip()
                        if text:
                            content_box = slide.shapes.add_textbox(
                                Inches(1), Inches(content_y), Inches(8), Inches(0.5)
                            )
                            content_frame = content_box.text_frame
                            content_frame.text = f"• {text}"
                            content_y += 0.6
                elif elem.get_text().strip():
                    text = elem.get_text().strip()
                    if len(text) > 10:  # 过滤太短的文本
                        content_box = slide.shapes.add_textbox(
                            Inches(0.5), Inches(content_y), Inches(9), Inches(0.8)
                        )
                        content_frame = content_box.text_frame
                        content_frame.text = text
                        content_y += 0.9
                        
                        if content_y > 6.5:  # 避免内容超出页面
                            break
        
        # 保存PPT
        prs.save(output_file)
        print(f"✅ PPT已保存: {output_file}")
        return output_file
    
    def method2_screenshot_slides(self, html_file: str, output_file: str = None) -> str:
        """
        方案2：浏览器截图生成PPTX
        
        优点：完美还原HTML样式和图表
        缺点：需要浏览器环境，生成的是图片PPT
        """
        if selenium is None:
            raise ImportError("需要安装selenium: pip install selenium")
        
        print("🔄 使用方案2：浏览器截图转换...")
        
        if output_file is None:
            output_file = html_file.replace('.html', '_screenshot.pptx')
        
        # 配置浏览器
        options = Options()
        options.add_argument('--headless')  # 无头模式
        options.add_argument('--no-sandbox')
        options.add_argument('--disable-dev-shm-usage')
        options.add_argument('--window-size=1920,1080')
        
        try:
            driver = webdriver.Chrome(options=options)
        except Exception as e:
            print(f"❌ Chrome浏览器启动失败: {e}")
            print("请确保安装了Chrome浏览器和chromedriver")
            return None
        
        try:
            # 打开HTML文件
            file_url = f"file://{os.path.abspath(html_file)}"
            driver.get(file_url)
            time.sleep(3)  # 等待页面加载
            
            # 创建PPT
            prs = Presentation()
            
            # 查找幻灯片数量
            slides = driver.find_elements(By.CSS_SELECTOR, '.slide, .ppt-slide, section')
            
            if not slides:
                # 如果没有找到幻灯片，截取整个页面
                slides = [driver.find_element(By.TAG_NAME, 'body')]
            
            for i, slide_elem in enumerate(slides):
                print(f"  截图第 {i+1} 页...")
                
                # 滚动到当前幻灯片
                driver.execute_script("arguments[0].scrollIntoView();", slide_elem)
                time.sleep(1)
                
                # 截图
                screenshot_path = f"temp_slide_{i}.png"
                slide_elem.screenshot(screenshot_path)
                
                # 添加到PPT
                slide_layout = prs.slide_layouts[5]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                
                # 添加图片
                slide.shapes.add_picture(
                    screenshot_path, 
                    Inches(0), Inches(0),
                    width=Inches(10), height=Inches(7.5)
                )
                
                # 删除临时文件
                os.remove(screenshot_path)
            
            # 保存PPT
            prs.save(output_file)
            print(f"✅ PPT已保存: {output_file}")
            return output_file
            
        finally:
            driver.quit()
    
    def method3_hybrid_approach(self, html_file: str, output_file: str = None) -> str:
        """
        方案3：混合方案
        
        结合内容解析和截图，生成更好的PPT
        """
        print("🔄 使用方案3：混合方案转换...")
        
        if output_file is None:
            output_file = html_file.replace('.html', '_hybrid.pptx')
        
        # 先用方案1解析内容
        temp_content_ppt = self.method1_parse_html_content(html_file, 
                                                          output_file.replace('.pptx', '_content.pptx'))
        
        # 如果selenium可用，为图表页面生成截图
        if selenium is not None:
            try:
                self.method2_screenshot_slides(html_file, 
                                             output_file.replace('.pptx', '_screenshot.pptx'))
                print("✅ 混合方案完成，生成了内容版和截图版两个PPT文件")
            except Exception as e:
                print(f"⚠️  截图功能失败: {e}")
        
        return temp_content_ppt
    
    async def convert_html_to_ppt(self, html_file: str, method: str = "parse", 
                                  output_file: str = None) -> str:
        """
        转换HTML为PPT的主接口
        
        Args:
            html_file: HTML文件路径
            method: 转换方法 ("parse", "screenshot", "hybrid")
            output_file: 输出PPT文件路径
        
        Returns:
            生成的PPT文件路径
        """
        if not os.path.exists(html_file):
            raise FileNotFoundError(f"HTML文件不存在: {html_file}")
        
        print(f"🚀 开始转换 {html_file} -> PPT")
        
        if method == "parse":
            return self.method1_parse_html_content(html_file, output_file)
        elif method == "screenshot":
            return self.method2_screenshot_slides(html_file, output_file)
        elif method == "hybrid":
            return self.method3_hybrid_approach(html_file, output_file)
        else:
            raise ValueError(f"不支持的转换方法: {method}")


class EnhancedPPTGenerator:
    """增强的PPT生成器，支持直接生成PPTX格式"""
    
    def __init__(self):
        """初始化增强PPT生成器"""
        from ppt_generator_demo import PPTGeneratorDemo
        self.html_generator = PPTGeneratorDemo()
        self.converter = HTMLToPPTConverter()
    
    async def generate_pptx_directly(self, task: str, file_names: List[str] = None,
                                   method: str = "parse", output_file: str = None) -> str:
        """
        直接生成PPTX格式的PPT
        
        Args:
            task: PPT生成任务描述
            file_names: 输入文件列表
            method: 转换方法
            output_file: 输出文件路径
        
        Returns:
            生成的PPTX文件路径
        """
        print("🎯 开始生成PPTX格式PPT...")
        
        # 1. 先生成HTML格式PPT
        print("📝 第1步：生成HTML格式PPT...")
        html_content = ""
        async for chunk in self.html_generator.generate_ppt(
            task=task, 
            file_names=file_names or [],
            stream=True
        ):
            html_content += chunk
        
        # 2. 保存HTML文件
        html_file = await self.html_generator.save_ppt_file(html_content, "temp_for_conversion.html")
        print(f"✅ HTML PPT已保存: {html_file}")
        
        # 3. 转换为PPTX
        print("🔄 第2步：转换为PPTX格式...")
        pptx_file = await self.converter.convert_html_to_ppt(html_file, method, output_file)
        
        # 4. 清理临时文件
        if os.path.exists(html_file):
            os.remove(html_file)
            print("🧹 临时HTML文件已清理")
        
        return pptx_file


async def demo_html_to_ppt_conversion():
    """演示HTML转PPT的各种方法"""
    print("=" * 60)
    print("🎯 HTML转PPT转换演示")
    print("=" * 60)
    
    # 检查是否有现成的HTML PPT文件
    output_dir = Path("output")
    html_files = list(output_dir.glob("*.html")) if output_dir.exists() else []
    
    if not html_files:
        print("📝 未找到现有的HTML PPT文件，先生成一个...")
        from ppt_generator_demo import PPTGeneratorDemo
        demo = PPTGeneratorDemo()
        await demo.demo_simple_ppt()
        html_files = list(output_dir.glob("*.html"))
    
    if html_files:
        html_file = str(html_files[0])
        print(f"📄 使用HTML文件: {html_file}")
        
        converter = HTMLToPPTConverter()
        
        print("\n1️⃣ 方案1：HTML内容解析转换")
        try:
            pptx_file1 = await converter.convert_html_to_ppt(html_file, "parse")
            print(f"✅ 生成成功: {pptx_file1}")
        except Exception as e:
            print(f"❌ 方案1失败: {e}")
        
        if selenium is not None:
            print("\n2️⃣ 方案2：浏览器截图转换")
            try:
                pptx_file2 = await converter.convert_html_to_ppt(html_file, "screenshot")
                print(f"✅ 生成成功: {pptx_file2}")
            except Exception as e:
                print(f"❌ 方案2失败: {e}")
        
        print("\n3️⃣ 方案3：混合方案转换")
        try:
            pptx_file3 = await converter.convert_html_to_ppt(html_file, "hybrid")
            print(f"✅ 生成成功: {pptx_file3}")
        except Exception as e:
            print(f"❌ 方案3失败: {e}")
    
    else:
        print("❌ 未找到HTML PPT文件进行转换")


async def demo_direct_pptx_generation():
    """演示直接生成PPTX格式PPT"""
    print("=" * 60)
    print("🎯 直接生成PPTX格式PPT演示")
    print("=" * 60)
    
    generator = EnhancedPPTGenerator()
    
    try:
        pptx_file = await generator.generate_pptx_directly(
            task="生成一份关于人工智能技术发展的PPT报告",
            method="parse",
            output_file="ai_tech_report.pptx"
        )
        print(f"🎉 PPTX格式PPT生成成功: {pptx_file}")
    except Exception as e:
        print(f"❌ 直接生成PPTX失败: {e}")


if __name__ == "__main__":
    print("🚀 启动HTML转PPT转换器...")
    
    # 检查依赖
    print("📋 检查依赖包...")
    missing_deps = []
    
    try:
        import pptx
        print("✅ python-pptx 已安装")
    except ImportError:
        missing_deps.append("python-pptx")
    
    try:
        import bs4
        print("✅ beautifulsoup4 已安装")
    except ImportError:
        missing_deps.append("beautifulsoup4")
    
    if selenium is None:
        print("⚠️  selenium 未安装（截图功能将不可用）")
    else:
        print("✅ selenium 已安装")
    
    if missing_deps:
        print(f"❌ 缺少依赖包: {', '.join(missing_deps)}")
        print(f"请安装: pip install {' '.join(missing_deps)}")
        sys.exit(1)
    
    # 运行演示
    import asyncio
    
    print("\n请选择演示模式：")
    print("1. HTML转PPT转换演示")
    print("2. 直接生成PPTX格式PPT演示")
    
    choice = input("请选择 (1 或 2): ").strip()
    
    try:
        if choice == "1":
            asyncio.run(demo_html_to_ppt_conversion())
        elif choice == "2":
            asyncio.run(demo_direct_pptx_generation())
        else:
            print("❌ 无效选择")
    except KeyboardInterrupt:
        print("\n⏹️  用户中断操作")
    except Exception as e:
        print(f"❌ 运行出错: {e}")
